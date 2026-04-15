import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def convert(input_path: str, output_path: str, progress_cb=None):
    doc = Document(input_path)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # 按 Heading 1 分组
    slides_data = []
    current = {"title": "", "content": [], "images": []}

    for para in doc.paragraphs:
        style = para.style.name
        if style == "Heading 1":
            if current["title"] or current["content"]:
                slides_data.append(current)
            current = {"title": para.text, "content": [], "images": []}
        elif style == "Heading 2":
            current["content"].append(("h2", para.text))
        else:
            if para.text.strip():
                current["content"].append(("body", para.text))
        for run in para.runs:
            for elem in run._r:
                if elem.tag.endswith("}drawing") or elem.tag.endswith("}pict"):
                    current["images"].append(None)

    if current["title"] or current["content"]:
        slides_data.append(current)

    total = len(slides_data) or 1
    for idx, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        # 标题
        if sd["title"]:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd["title"]
            p.font.size = Pt(36)
            p.font.bold = True

        # 内容
        if sd["content"]:
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.4))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, (kind, text) in enumerate(sd["content"]):
                p = tf2.paragraphs[i] if i == 0 else tf2.add_paragraph()
                p.text = text
                if kind == "h2":
                    p.font.size = Pt(24)
                    p.font.bold = True
                else:
                    p.font.size = Pt(18)

        if progress_cb:
            progress_cb(int((idx + 1) / total * 100), f"第 {idx+1}/{total} 张")

    prs.save(output_path)
    if progress_cb:
        progress_cb(100)

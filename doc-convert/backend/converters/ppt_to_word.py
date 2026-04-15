import os
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DInches
import io


def convert(input_path: str, output_path: str, progress_cb=None):
    prs = Presentation(input_path)
    doc = Document()
    total = len(prs.slides)

    for idx, slide in enumerate(prs.slides):
        title_text = ""
        content_texts = []
        images = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                if shape.shape_type == 13 or (hasattr(shape, "placeholder_format") and
                        shape.placeholder_format and shape.placeholder_format.idx == 0):
                    title_text = shape.text_frame.text.strip()
                else:
                    content_texts.extend(texts)
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_blob = shape.image.blob
                    images.append(image_blob)
                except Exception:
                    pass

        if title_text:
            doc.add_heading(title_text, level=1)
        for t in content_texts:
            doc.add_paragraph(t)
        for blob in images:
            doc.add_picture(io.BytesIO(blob), width=DInches(4))

        if idx < total - 1:
            doc.add_page_break()

        if progress_cb:
            progress_cb(int((idx + 1) / total * 100), f"第 {idx+1}/{total} 张")

    doc.save(output_path)
    if progress_cb:
        progress_cb(100)

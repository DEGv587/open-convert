import os
import fitz
from pptx import Presentation
from pptx.util import Inches


def convert(input_path: str, output_path: str, progress_cb=None):
    doc = fitz.open(input_path)
    total = len(doc)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    work_dir = os.path.dirname(output_path)

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_path = os.path.join(work_dir, f"_slide_{i:04d}.png")
        pix.save(img_path)

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

        if progress_cb:
            progress_cb(int((i + 1) / total * 95), f"第 {i+1}/{total} 页")

    prs.save(output_path)

    # 清理临时图片
    for i in range(total):
        img_path = os.path.join(work_dir, f"_slide_{i:04d}.png")
        if os.path.exists(img_path):
            os.remove(img_path)

    if progress_cb:
        progress_cb(100)
